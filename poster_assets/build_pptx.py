"""Fill the Stanford CS poster template with content + figures -> editable PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CARD = RGBColor(0x8C, 0x15, 0x15); INK = RGBColor(0x1a, 0x1a, 0x1a); GREY = RGBColor(0x55, 0x55, 0x55)
A = "/Users/manrajmondair/bci-review/poster_assets"
TPL = "/Users/manrajmondair/Downloads/Manraj Singh Mondair CS 281.pptx"
OUT = "/Users/manrajmondair/Downloads/BCI_Poster_Mondair_EDITABLE.pptx"

prs = Presentation(TPL)
slide = prs.slides[0]
by_id = {sp.shape_id: sp for sp in slide.shapes}

def move(sp, l, t, w, h):
    sp.left, sp.top, sp.width, sp.height = Inches(l), Inches(t), Inches(w), Inches(h)

def fill(sp, items, size=14, color=INK, bullet=True, lead=1.0, space_after=6):
    tf = sp.text_frame; tf.word_wrap = True
    tf.clear()
    for i, it in enumerate(items):
        txt = it; bold = False
        if isinstance(it, tuple): txt, bold = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after); p.line_spacing = lead
        r = p.add_run(); r.text = (("•  ") if bullet else "") + txt
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Lato"; r.font.bold = bold

def add_text(l, t, w, h, runs, size, color=INK, bold=False, italic=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = runs
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.italic = italic; r.font.name = "Lato"
    return tb

def add_img(path, l, t, w):
    return slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))

# ---------------- LEFT: Project Overview (id 140) ----------------
po = by_id[140]; move(po, 0.5, 4.9, 7.85, 10.2)
fill(po, [
 "BCIs decode EEG into commands for assistive devices and consumer neurotech. We ask the privacy question task benchmarks ignore: does the signature that enables decoding also reveal WHO the user is?",
 "Motor-imagery BCI models leak stable subject identity as a side-channel — recoverable by an attacker holding only model outputs or embeddings.",
 "Headline: 100% subject re-identification across 104 people (chance 0.96%) — while the decoder reaches only ~35% task accuracy. Privacy leakage is decoupled from task utility.",
 "Ethics: a model trained to read commands can fingerprint its cohort. EEG should be treated as biometric data under GDPR Article 9 and emerging neurorights.",
 "Contribution: an audit-clean benchmark — 3 corpora x 5 attacks x 4 defenses x adaptive attackers, every number with bootstrap CIs + a 240-invariant audit.",
], size=13.5)
if 145 in by_id: by_id[145].text_frame.clear()

# ---------------- LEFT: Datasets & Metrics (table + id 146) ----------------
rows = [("Corpus","Subj.","Role"),("PhysioNet EEG-MMIDB","104","primary"),
        ("BCI IV-2a","9x2","cross-session"),("Lee 2019 OpenBMI","54x2","replication")]
tbl_sh = slide.shapes.add_table(4, 3, Inches(0.45), Inches(16.0), Inches(7.85), Inches(2.3)).table
tbl_sh.columns[0].width=Inches(4.6); tbl_sh.columns[1].width=Inches(1.2); tbl_sh.columns[2].width=Inches(2.05)
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        c = tbl_sh.cell(ri,ci); c.text = val
        pr = c.text_frame.paragraphs[0]; rn = pr.runs[0]
        rn.font.size = Pt(13); rn.font.name="Lato"
        if ri==0: rn.font.bold=True; rn.font.color.rgb=RGBColor(0xff,0xff,0xff); c.fill.solid(); c.fill.fore_color.rgb=CARD
        else: rn.font.color.rgb=INK
met = by_id[146]; move(met, 0.45, 18.6, 7.85, 4.0)
fill(met, [
 "Closed-set top-1 — name the subject among N enrolled (chance 1/N).",
 "Open-set AUC / EER — same-vs-different verification on subjects NEVER seen in training (chance 0.5).",
 "Membership-inference AUC + advantage — was this person in the training set?",
 "Utility cost — motor-imagery accuracy (pp) a defense sacrifices.",
], size=12)

# ---------------- MIDDLE: Methods (framing + schematic + id 148 bullets) ----------------
add_text(9.05, 4.85, 14.4, 0.9, "Three victim decoders, attacked five ways (A1-A5); four defenses (D1-D4) stress-tested against generic AND adaptive attackers.",
         size=13, color=CARD, italic=True)
add_img(f"{A}/fig_schematic.png", 9.05, 5.7, 14.4)
me = by_id[148]; move(me, 9.05, 12.0, 14.4, 3.8)
fill(me, [
 "Victims: FBCSP+LDA and Riemannian tangent-space+LR (classical); EEGNet (deep CNN); contrastive EEGNet embedder for open-set verification.",
 "Adaptive threat model — the real test: the attacker knows the defense and re-trains the encoder end-to-end on identity labels (encoder fine-tune).",
 "Rigor: 1000-resample trial-grouped bootstrap CIs; shuffled-label negative control; run provenance + a 240-invariant audit on every commit.",
], size=12.5)

# ---------------- MIDDLE: Discussions / Future / References (id 152) ----------------
dis = by_id[152]; move(dis, 9.05, 16.85, 14.4, 6.6)
tf = dis.text_frame; tf.word_wrap = True; tf.clear()
def head_para(p, text):
    r=p.add_run(); r.text=text; r.font.size=Pt(13.5); r.font.bold=True; r.font.color.rgb=CARD; r.font.name="Lato"
def bul(p, text, size=12):
    r=p.add_run(); r.text="•  "+text; r.font.size=Pt(size); r.font.color.rgb=INK; r.font.name="Lato"
p=tf.paragraphs[0]; head_para(p,"Discussions:")
for t in [
 "Identity leaks through every decoder family; survives task change, session change, even resting state — and verifies subjects the model never saw (AUC 0.92 on two corpora).",
 "Ad-hoc transforms and DANN cut leakage vs a generic probe but COLLAPSE under an adaptive attacker (re-ID restored to 0.80) — security by obscurity fails.",
 "Only DP-SGD is adaptively robust: e<=1 blocks both re-ID and DP-aware membership inference at ~6pp task cost; e=3 stops re-ID but NOT membership inference."]:
    p=tf.add_paragraph(); p.space_after=Pt(5); p.line_spacing=1.0; bul(p,t)
p=tf.add_paragraph(); p.space_before=Pt(4); head_para(p,"Future Research:")
p=tf.add_paragraph(); p.space_after=Pt(5); bul(p,"Population-scale cohorts (TUH-EEG ~10^4), other paradigms (sleep, ERP), tighter federated budgets, per-subject ancestry metadata for fairness audits.")
p=tf.add_paragraph(); p.space_before=Pt(4)
r=p.add_run(); r.text="References:  "; r.font.size=Pt(10); r.font.bold=True; r.font.color.rgb=CARD; r.font.name="Lato"
for ref in ["[1] Schalk+ BCI2000/PhysioNet 2004.  [2] Lee+ OpenBMI 2019.  [3] Lawhern+ EEGNet 2018.  [4] Abadi+ DP-SGD 2016.  [5] Shokri+ MIA 2017.  [6] Yeom+ 2018."]:
    rr=p.add_run(); rr.text=ref; rr.font.size=Pt(10); rr.font.color.rgb=GREY; rr.font.name="Lato"
if 147 in by_id: by_id[147].text_frame.clear()

# ---------------- RIGHT: Results sub-headers + figures + callout ----------------
sh1 = by_id[149]; move(sh1, 24.3, 4.85, 11.3, 0.6)
sh1.text_frame.clear(); r=sh1.text_frame.paragraphs[0].add_run()
r.text="Attacks — identity leaks at chance task accuracy"; r.font.size=Pt(18); r.font.bold=True; r.font.color.rgb=CARD; r.font.name="Lato"
add_text(24.3, 5.5, 11.2, 0.8, "Classical decoders re-identify all 104 subjects perfectly; the signal survives task/session change and generalizes to unseen people.", size=12, color=INK)
add_img(f"{A}/fig_closed_set.png", 24.6, 6.35, 10.6)
add_text(24.3, 10.35, 11.2, 0.7, "100% re-ID (Riemann) at 0.96% chance — identity is trivially recoverable from a model trained only to decode commands.", size=11.5, color=GREY)
# callout box
co = slide.shapes.add_textbox(Inches(24.3), Inches(11.2), Inches(11.2), Inches(1.7))
co.fill.solid(); co.fill.fore_color.rgb=RGBColor(0xF7,0xEC,0xEC); co.line.color.rgb=CARD; co.line.width=Pt(2)
ctf=co.text_frame; ctf.word_wrap=True
p=ctf.paragraphs[0]; r=p.add_run(); r.text="AUC 0.925 / 0.920"; r.font.size=Pt(30); r.font.bold=True; r.font.color.rgb=CARD; r.font.name="Lato"
p2=ctf.add_paragraph(); r=p2.add_run(); r.text="Open-set verification on subjects NEVER seen in training — PhysioNet 0.925, replicates 0.920 on Lee 2019. EEG behaves like a biometric template."
r.font.size=Pt(12); r.font.color.rgb=INK; r.font.name="Lato"

sh2 = by_id[151]; move(sh2, 24.3, 13.1, 11.3, 0.6)
sh2.text_frame.clear(); r=sh2.text_frame.paragraphs[0].add_run()
r.text="Defenses — only DP-SGD survives adaptive attack"; r.font.size=Pt(18); r.font.bold=True; r.font.color.rgb=CARD; r.font.name="Lato"
add_img(f"{A}/fig_dp_sweep.png", 25.4, 13.75, 9.0)
add_text(24.3, 18.1, 11.2, 0.7, "DP-SGD privacy-utility frontier: e<=1 holds the adaptive encoder-fine-tune attacker near chance at ~6pp task cost.", size=11.5, color=GREY)
add_img(f"{A}/fig_dpmia.png", 26.2, 18.85, 7.3)
add_text(24.3, 22.4, 11.2, 0.7, "DP-aware membership inference: e=3 fails (AUC 0.89 ~ undefended); only e<=1 reaches chance, tracking the Yeom (2018) bound.", size=11.5, color=GREY)

prs.save(OUT)
print("saved", OUT)
